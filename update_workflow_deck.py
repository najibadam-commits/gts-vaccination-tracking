"""Update the GTS Pipeline Process Workflow deck to match the current pipeline.

Edits the supplied deck in place rather than rebuilding it: every paragraph,
table row and new slide is cloned from one already in the file, so the
deck's own fonts, colours, spacing and layout survive untouched. Nothing here
sets a font size or a colour directly.
"""
import copy
import sys

from pptx import Presentation
from pptx.text.text import _Paragraph

SRC = sys.argv[1]
DST = sys.argv[2]


# --------------------------------------------------------------- helpers
def set_para(para, text):
    """Replace a paragraph's text, keeping the first run's formatting."""
    runs = para.runs
    if not runs:
        para.add_run().text = text
        return para
    runs[0].text = text
    for r in runs[1:]:
        r._r.getparent().remove(r._r)
    return para


def paras(shape):
    return list(shape.text_frame.paragraphs)


def find_shape(slide, startswith=None, name=None):
    for sh in slide.shapes:
        if name and sh.name == name:
            return sh
        if startswith and sh.has_text_frame and \
                sh.text_frame.text.strip().startswith(startswith):
            return sh
    return None


A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def space_after(para, points):
    """Set a paragraph's space-after, in points x100 (OOXML spcPts units).

    The deck gives each metric's description a 7pt gap to the next metric, and
    zeroes it on the last one in the column. Inserting a metric therefore has
    to move that zero along, or the new last pair sits tight against the one
    above it while the old last pair keeps the gap.
    """
    pPr = para._p.get_or_add_pPr()
    for old in pPr.findall(f"{A_NS}spcAft"):
        pPr.remove(old)
    spc_aft = pPr.makeelement(f"{A_NS}spcAft", {})
    pts = spc_aft.makeelement(f"{A_NS}spcPts", {"val": str(points)})
    spc_aft.append(pts)
    # spcAft must follow lnSpc/spcBef in the CT_TextParagraphProperties order
    ln_spc = pPr.find(f"{A_NS}lnSpc")
    (ln_spc if ln_spc is not None else pPr).addnext(spc_aft) \
        if ln_spc is not None else pPr.insert(0, spc_aft)


def clone_para_after(tf, template_para, text, after_para=None):
    """Add a paragraph styled like `template_para`, after `after_para`."""
    new_p = copy.deepcopy(template_para._p)
    anchor = (after_para or tf.paragraphs[-1])._p
    anchor.addnext(new_p)
    return set_para(_Paragraph(new_p, tf), text)


def set_row(table, idx, values):
    """Rewrite a table row's cells, keeping each cell's own run formatting."""
    for c, value in enumerate(values):
        if value is None:
            continue
        cell = table.cell(idx, c)
        set_para(cell.text_frame.paragraphs[0], value)
        for extra in list(cell.text_frame.paragraphs[1:]):
            extra._p.getparent().remove(extra._p)


def clone_row(table, template_idx, values):
    """Append a row cloned from an existing one, then fill it."""
    tbl = table._tbl
    new_tr = copy.deepcopy(tbl.tr_lst[template_idx])
    tbl.append(new_tr)
    set_row(table, len(table.rows) - 1, values)
    return len(table.rows) - 1


def scale_rows(table, total_emu):
    """Redistribute a table's height evenly over its (now more numerous) rows."""
    h = int(total_emu / len(table.rows))
    for r in table.rows:
        r.height = h


def duplicate_slide(prs, index, insert_after):
    """Copy a slide whole — layout, shapes, formatting — and place it.

    python-pptx has no slide-copy, and rebuilding the slide by hand would mean
    restating every font and colour the deck already defines. Deep-copying the
    shape tree onto a slide of the same layout keeps all of it.
    """
    src = prs.slides[index]
    dest = prs.slides.add_slide(src.slide_layout)
    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)
    for shp in src.shapes:
        dest.shapes._spTree.append(copy.deepcopy(shp._element))

    sld_id_lst = prs.slides._sldIdLst
    ids = list(sld_id_lst)
    moved = ids[-1]
    sld_id_lst.remove(moved)
    sld_id_lst.insert(insert_after + 1, moved)
    return dest


prs = Presentation(SRC)
S = list(prs.slides)

# ------------------------------------------------- slide 2 · exec summary
body = find_shape(S[1], startswith="The GTS pipeline converts")
for p in paras(body):
    if p.text.startswith("Coverage, time spent and the daily"):
        set_para(p, "Daily scope, cumulative scope, time spent and coverage each "
                    "follow one rule, applied once at source.")

# ---------------------------------------------- slide 4 · stage inventory
tbl = next(sh for sh in S[3].shapes if sh.has_table).table
set_row(tbl, 2, [None, None,
                 "Grid intersection → daily, cumulative and track-evidence status", None])
set_row(tbl, 3, [None, None,
                 "Team deployment, day-scoped time spent, ERM tabs", None])

# ------------------------------------- slide 5 · stage 2 visitation steps
col2 = find_shape(S[4], startswith="2 · Visitation Analysis")
for p in paras(col2):
    if p.text.startswith("Record day-only"):
        set_para(p, "Record day-only and carried-forward status in separate columns")
    elif p.text.startswith("Derive the daily and cumulative"):
        set_para(p, "Carry Track Evidence forward — tracks on any day count cumulatively")

# ------------------------------------------- slide 6 · stage 3 workbook
col1 = find_shape(S[5], startswith="3 · ERM Workbook")
for p in paras(col1):
    if p.text.startswith("Time Spent Analysis:"):
        set_para(p, "Time Spent: 08:00–15:00, scoped to the tracks' transmission date")
    elif p.text.startswith("Team time-range and target-population"):
        set_para(p, "Team time-range, efficiency, stationary and target-population tabs")

# ------------------------------------------------- slide 7 · charts & reporting
col1 = find_shape(S[6], startswith="4b · Report Charts")
for p in paras(col1):
    if p.text.startswith("Coverage donuts —"):
        set_para(p, "Coverage donuts — daily scoped to the day's planned settlements")
    elif p.text.startswith("Target population and household"):
        set_para(p, "Target-population donuts; efficiency scatter with a label gutter")

col2 = find_shape(S[6], startswith="5–7 · Reporting")
for p in paras(col2):
    if p.text.startswith("Daily: populate the organisation"):
        set_para(p, "Daily: fill the org PPTX — narration auto-fitted, overlap checked")

# ------------------------------------------------ slide 9 · error handling
tbl = next(sh for sh in S[8].shapes if sh.has_table).table
orig_h = sum(r.height for r in tbl.rows)
clone_row(tbl, len(tbl.rows) - 1, [
    "No missing-households column",
    "The indicator is left out of the narration, tables and workbook — never "
    "estimated to fill the gap"])
clone_row(tbl, len(tbl.rows) - 1, [
    "No day-N planning flag",
    "Daily coverage falls back to the full planned list, and the run log says so"])
scale_rows(tbl, orig_h)

# ----------------------------------------------- slide 10 · limitations
tbl = next(sh for sh in S[9].shapes if sh.has_table).table
orig_h = sum(r.height for r in tbl.rows)
clone_row(tbl, len(tbl.rows) - 1, [
    "Stationary flag is a prompt",
    "A team whose only settlement is genuinely tiny can land in the list; it is "
    "raised for supervisory verification, never asserted as a finding"])
scale_rows(tbl, orig_h)

# ------------------------------------------------------ slide 11 · KPIs
outcomes = find_shape(S[10], startswith="CAMPAIGN OUTCOMES")
ps = paras(outcomes)
for p in ps:
    if p.text.startswith("Follow-up list"):
        set_para(p, "Follow-up list by LGA, from assigned-settlement totals")
label_tpl = next(p for p in ps if p.text == "Teams under 12 minutes")
value_tpl = next(p for p in ps if p.text.startswith("Follow-up list"))
last = paras(outcomes)[-1]
lab = clone_para_after(outcomes.text_frame, label_tpl, "Stationary teams", after_para=last)
new_last = clone_para_after(outcomes.text_frame, value_tpl,
                            "Reported working, tracks show one location", after_para=lab)
space_after(last, 700)      # was the final metric, now has one below it
space_after(new_last, 0)    # is the final metric now

# --------------------------------------------- slide 12 · mermaid source
left = find_shape(S[11], startswith="flowchart LR")
for p in paras(left):
    if "team_time_range" in p.text:
        set_para(p, "    AN([team_time_range / team_efficiency / target_population])")

# ============================================ NEW SLIDE · daily vs cumulative
# Cloned from the "Stage Detail — Charts & Reporting" slide (index 6) so the
# two-column styling, title treatment and label formatting are the deck's own.
new = duplicate_slide(prs, 6, insert_after=6)
# This template's title on the two-column layout is a body placeholder named
# "Title 3", so slide.shapes.title is None — it is found by its text instead.
title_shape = find_shape(new, startswith="Stage Detail")
set_para(title_shape.text_frame.paragraphs[0], "Daily vs Cumulative — the Rule")

# Each column is written into the source column's paragraphs IN ORDER, so the
# rebuilt slide keeps the original's rhythm: heading, sub-heading, description,
# then LABEL / value, LABEL / bullets, LABEL / value, LABEL / value. Writing by
# category instead of by position is what put "WHY IT MATTERS" after its own
# bullets on the first attempt.
LEFT = {
    "heading": "DAILY",
    "sub": "the reporting day alone",
    "desc": "Both sides of the ratio are scoped to the day.",
    "label_1": "SCOPE",
    "value_1": "Denominator: settlements the day1 / day2 planning flag says were "
               "scheduled. Numerator: day_{N}_daily, from this run's tracks only.",
    "label_2": "WHY IT MATTERS",
    "bullets": [
        "Against the whole campaign list, a normal day reads as a failure",
        "Donut, LGA bar and narration all apply the same rule, so they agree",
        "Time spent keeps only the latest transmission date in the export",
        "A team that reported yesterday but not today shows as no tracks",
    ],
    "label_3": "WATCH FOR",
    "value_2": "No flag for that day → the full planned list is used, and the run "
               "log says so.",
    "label_4": "DRIVES",
    "value_3": "State Daily Coverage · Time Spent Analysis",
}
RIGHT = {
    "heading": "CUMULATIVE",
    "sub": "Day 1 through the reporting day",
    "desc": "Evidence of tracks on any day counts, whatever day was planned.",
    "label_1": "SOURCE",
    "value_1": "Track Evidence carried forward from the previous day, ORed with "
               "today's — Grid Cells Visited > 0, or track_count > 0.",
    "label_2": "WHY IT MATTERS",
    "bullets": [
        "Applied once in stage 2, so every consumer agrees on what is reached",
        "A settlement worked on Day 1 keeps counting on Day 3",
        "Settlements with pings but no gridded polygon stop reading as unvisited",
        "Post-campaign time spent spans the whole period on purpose",
    ],
    "label_3": "WATCH FOR",
    "value_2": "Totals read higher than before for settlements with pings but no "
               "polygon in the gridded layer.",
    "label_4": "DRIVES",
    "value_3": "State Cumulative Coverage · Cumulative Settlements Coverage",
}


def fill_column(shape, spec):
    """Rewrite a stage-detail column, reusing each paragraph's style in place.

    The source column's paragraph sequence is fixed by the deck's own design:

        heading · sub · description · LABEL · value · LABEL · bullet×4 ·
        LABEL · value · LABEL · value

    Each new line is written into the paragraph already carrying the style it
    needs, in that order, so nothing has to be restyled and nothing moves.
    """
    tf = shape.text_frame
    ps = paras(shape)
    set_para(ps[0], spec["heading"])
    set_para(ps[1], spec["sub"])
    set_para(ps[2], spec["desc"])

    order = ["label_1", "value_1", "label_2", "bullets", "label_3", "value_2",
             "label_4", "value_3"]
    bullets = [p for p in ps[3:] if p.level == 1]
    non_bullets = [p for p in ps[3:] if p.level != 1]

    nb = iter(non_bullets)
    for key in order:
        if key == "bullets":
            lines = spec["bullets"]
            for i, line in enumerate(lines):
                if i < len(bullets):
                    set_para(bullets[i], line)
                else:
                    clone_para_after(tf, bullets[0], line, after_para=bullets[-1])
            for extra in bullets[len(lines):]:
                extra._p.getparent().remove(extra._p)
        else:
            set_para(next(nb), spec[key])
    for extra in nb:
        extra._p.getparent().remove(extra._p)


cols = [sh for sh in new.shapes
        if sh.has_text_frame and sh is not title_shape and sh.top > title_shape.top]
cols.sort(key=lambda sh: sh.left)
fill_column(cols[0], LEFT)
fill_column(cols[1], RIGHT)

prs.save(DST)
print(f"saved {DST} — {len(prs.slides)} slides")
