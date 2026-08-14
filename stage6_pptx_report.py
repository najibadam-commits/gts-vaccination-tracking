"""Stage 6 — Daily Tracking Report in the organization's PPTX template.

Populates the eHealth Africa "GTS Tracking Report" deck (Google-Slides
export) with this campaign day's numbers, charts and maps. Slide numbers
below are the FINAL output deck's order (what you'd count opening the
generated .pptx) — not the same as the SLIDE_* constants further down,
which name the source slideN.xml files in the original template and don't
shift when slides are inserted/removed/merged at build time:
  slide 1  title (state, day, date)
  slide 2  background (LGA/ward/team/settlement counts + implementation map)
  slide 3  team deployed / reported (+ deployment bar)
  slide 4  time spent analysis (+ time-spent bar) — the template has no
           dedicated slide for this, so one is created at build time by
           duplicating the Team Deployment slide's layout (title +
           narrative + one chart) and inserting it right after slide 3
  slide 5  state DAILY coverage — today's planned subset only (donut + LGA bar)
  slide 6  cumulative coverage — donut + LGA breakdown bar, side by side
           (merged from the template's separate "cumulative settlements
           coverage" and "cumulative coverage by LGA" slides; laid out with
           the same two-chart, side-by-side design as the daily coverage
           slide, reusing that slide's own picture boxes so the two stay
           visually consistent)
  slide 7  [removed — day-over-day settlement table not auto-filled yet]
  slide 8  statewide coverage map
  slide 9+ one slide per implementing LGA, coverage map
  challenges / photos slides — left as editable placeholders
  thank-you slide — unchanged

Requires the template's slide layout to match the one this module was
built against (see gts_pipeline/README.md). If the org updates the
template, re-run the slide-inspection steps described there before
changing the constants below.
"""
import argparse
import copy
import math
import os
import re
import shutil
import sys
import tempfile
from datetime import date

import pandas as pd
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.text.text import _Paragraph
from pptx.util import Emu, Pt

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from pptx_xml_utils import safe_extract, rezip, duplicate_slide, delete_slide, clean_unused_slides
from stage2_analysis import day_scope_mask, detect_day_column, find_col

# ------------------------------------------------------------------ geometry
# Slide is 10 x 5.625in. Everything below is expressed in inches and converted
# once, so the numbers stay readable and can be checked against a slide ruler.
def _in(inches: float) -> Emu:
    return Emu(int(round(inches * 914400)))


SLIDE_W_IN, SLIDE_H_IN = 10.0, 5.625

# ------------------------------------------------------------- template map
# Slide numbers as they appear in the template's ppt/slides/slideN.xml
# (verified to match presentation order 1:1 for this specific template).
SLIDE_TITLE = "slide1.xml"
SLIDE_BACKGROUND = "slide2.xml"
SLIDE_DEPLOY = "slide3.xml"
SLIDE_DAILY_COVERAGE = "slide4.xml"
SLIDE_CUMULATIVE_DONUT = "slide5.xml"
SLIDE_CUMULATIVE_LGA = "slide6.xml"        # removed — its chart merges onto SLIDE_CUMULATIVE_DONUT
SLIDE_DAY_COMPARE_TABLE = "slide7.xml"     # removed — not auto-filled
SLIDE_STATE_MAP = "slide8.xml"
SLIDE_LGA_MAP_TEMPLATE = "slide9.xml"       # kept as LGA slide #1, duplicated for the rest
SLIDE_LGA_MAP_ORIGINALS_TO_DROP = [f"slide{n}.xml" for n in range(10, 22)]  # 10..21

# Slides inside the drop range above that must survive regardless. The range was
# written on the assumption that slides 10-21 are all hand-made LGA map slides,
# but the template also carries a Challenges slide in there, which the blanket
# delete was removing. Matching on slide text rather than a slide number keeps
# this working if the template is re-ordered.
PRESERVE_SLIDE_KEYWORDS = ("challenge",)

# A generous, consistent map box (EMU) reused for the statewide map and every
# LGA map slide — the template's per-LGA slide (slide9) was hand-sized for
# one specific LGA's screenshot and is too small/narrow for the rest.
#
# Its 1.414:1 aspect is the A-series page ratio, so both the A3 single-panel
# maps and the A0 side-by-side pages fill it edge to edge with no letterboxing.
# It is also close to the largest box this 10 x 5.625in slide can give a map of
# that shape once the title bar is accounted for, so there is little to gain by
# enlarging it for the wider two-panel maps.
#
# Top moved to 1.10in so it clears the title placeholder (which runs to 1.06in)
# rather than tucking under its lower edge, and the height cut so the lower
# edge stops above the footer band at 5.25in. The 1.414:1 A-series ratio is
# preserved, so both the A3 single-panel maps and the A0 side-by-side pages
# still fill the box with no letterboxing.
MAP_BOX = (_in(2.12), _in(1.10), _in(5.76), _in(4.07))

# ---- narration band --------------------------------------------------------
# Every generated slide puts its narration in a FIXED band between the title
# and whatever sits below it, with the font fitted to that band. The template's
# own narration box is 10in wide at 18pt with no height limit, which is exactly
# what let a two-sentence narration grow downwards into the charts. Pinning the
# band and shrinking the type to fit removes the failure mode rather than
# tuning each slide's wording and hoping it still fits next campaign.
NARR_LEFT_IN = 0.25
NARR_WIDTH_IN = 9.50
# Clear of the title placeholder, whose box runs to 1.06in. The template's own
# narration boxes start at 0.90 and technically sit inside it; harmless with
# the template's short titles, but this pipeline's titles are longer and the
# margin is not worth the risk.
NARR_TOP_IN = 1.08
NARR_GAP_IN = 0.12          # clear space between narration and the content below
NARR_BASE_PT = 14.0
NARR_MIN_PT = 9.0
# Average glyph width as a fraction of font size for the template's body face,
# and line height as a multiple of it. Deliberately pessimistic: over-estimating
# only shrinks the text a little, under-estimating puts it back in the chart.
#
# Both were raised after a real run still showed the Team Time Efficiency
# narration touching its table. The deck is a Google Slides export whose font is
# substituted at render time, and the substitute is wider and more generously
# leaded than the original — an estimate tuned to the original under-counts the
# lines. These figures are set for the substitute.
NARR_CHAR_W_RATIO = 0.56
NARR_LINE_H_RATIO = 1.38

# ---- content band ----------------------------------------------------------
# Top of the charts/tables on the slides this module builds, and the bottom
# limit they must respect. A common top is what makes the deck read as one
# layout rather than a set of slides each nudged into place separately.
CONTENT_TOP_IN = 1.95
# The master's footer band starts at 5.25in. Nothing may cross it — a chart or
# map that does prints over the eHealth Africa URL and copyright line, which is
# how the map slides ended up with their lower edge inside the green bar.
CONTENT_BOTTOM_IN = 5.18

# The two coverage slides (daily, cumulative) carry two charts side by side.
# The template's own picture boxes are unequal — one starts at 1.66in and the
# other at 1.95, with different heights — which made the two charts sit at
# visibly different heights and put the taller one under the title. One
# explicit pair is used for both slides instead, so they can be read against
# each other and neither reaches the narration.
COVERAGE_TOP_IN = 1.72
COVERAGE_PIC_BOXES = [
    (_in(0.15), _in(COVERAGE_TOP_IN), _in(4.70), _in(3.44)),
    (_in(4.98), _in(COVERAGE_TOP_IN), _in(4.87), _in(3.44)),
]

# Time Spent Analysis: time-range chart left, under-12-minutes follow-up right.
TIME_CHART_BOX = (_in(0.20), _in(CONTENT_TOP_IN), _in(5.45), _in(3.22))
TIME_TABLE_BOX = (_in(5.78), _in(CONTENT_TOP_IN), _in(4.02), _in(2.30))
TIME_TABLE_MAX_ROWS = 9
TIME_TABLE_COL_WIDTHS = (2.42, 1.60)          # LGA, Teams <12 min
TIME_TABLE_ROW_HEIGHT = _in(0.17)
TIME_TABLE_FONT_PT = 8
TIME_TABLE_CAPTION = "Top LGAs by teams spending under 12 minutes in the field."

# PowerPoint and LibreOffice both treat a row's height as a MINIMUM and grow it
# to fit the text plus the cell's own margins, so the requested height above is
# a floor, not the rendered height. Cell margins are cut right down (below) to
# let rows get close to it, and this is what an 8pt row actually comes out at —
# used to place the caption, since the shape's reported height is the requested
# one and would put the caption underneath the table's last rows.
TABLE_RENDERED_ROW_IN = 0.185
TIME_TABLE_CELL_MARGIN_V = Emu(9144)    # 0.01in top/bottom
TIME_TABLE_CELL_MARGIN_H = Emu(45720)   # 0.05in left/right
TABLE_CAPTION_PT = 8

# Target Population & Household Coverage. Panels sit in the content band; the
# narration above them is fitted to the band, so it can no longer reach them.
TARGET_SLIDE_TITLE = "Target Population Coverage Analysis"
TARGET_BOXES_3 = [
    (_in(0.18), _in(CONTENT_TOP_IN), _in(2.60), _in(3.22)),
    (_in(2.88), _in(CONTENT_TOP_IN), _in(2.60), _in(3.22)),
    (_in(5.58), _in(CONTENT_TOP_IN), _in(4.24), _in(3.22)),
]
# Used when the settlement list carries only one of the two estimates, so the
# remaining donut and the LGA bar share the width instead of leaving a gap.
TARGET_BOXES_2 = [
    (_in(0.40), _in(CONTENT_TOP_IN), _in(4.20), _in(3.22)),
    (_in(4.90), _in(CONTENT_TOP_IN), _in(4.70), _in(3.22)),
]

# Team Time Efficiency: the quadrant scatter, full width, and nothing else.
#
# The flagged-team table that used to sit to its right has been removed. It was
# the element that kept colliding with the narration, and the same teams are
# already named in the chart's own label gutter and listed in full in the
# workbook's "Teams Flagged (Time x Cov)" tab — so the slide loses no
# information by dropping it, and the chart gains the whole width, which is
# what makes those team codes legible at presentation size.
EFF_SLIDE_TITLE = "Team Time Efficiency"
# This slide's narration carries two findings (unproductive field time, and
# stationary teams), so it runs longer than the others and needs a deeper band.
EFF_CONTENT_TOP_IN = 2.05
EFF_CHART_LEFT_IN = 0.30
EFF_CHART_WIDTH_IN = 9.40
# Capped below NARR_BASE_PT so this slide's narration is compact by
# construction rather than only when the fitter happens to shrink it.
EFF_NARR_MAX_PT = 12.0
# Clear air between the last line of the narration and the top of the chart.
EFF_NARR_CLEARANCE_IN = 0.20

# Photo slides are retitled from the template's own wording. Keyed on the
# template text rather than a slide number, which shifts with the number of
# LGA map slides inserted ahead of them.
PICTURE_SLIDE_TITLES = {
    "Supportive Supervision Pictures": "ERM & Supportive Supervision Pictures",
}

# ---- title slide -----------------------------------------------------------
# The landing page carries three lines, in this order and no other:
#
#     {State} State                          bold
#     GTS Tracking Report                    bold
#     Evening Review Meeting Presentation
#
# Sized as a hierarchy rather than three equal lines: the middle line is the
# report's name and takes the weight, the first names the campaign's state and
# the third is a subtitle. Set at one size they would not fit the placeholder's
# width — the third line alone is 35 characters, which wraps at the template's
# 40pt — and would read as three competing headlines.
TITLE_LINE_1 = "{state} State"
TITLE_LINE_2 = "GTS Tracking Report"
TITLE_LINE_3 = "Evening Review Meeting Presentation"
TITLE_LINE_SIZES = (28.0, 40.0, 18.0)
TITLE_LINE_BOLD = (True, True, False)


# --------------------------------------------------------------- utilities
def _slide_text(unpacked_dir: str, slide_filename: str) -> str:
    """All rendered text on a slide, read straight from its unpacked XML.

    Used to identify slides by content before python-pptx opens the package,
    so a slide can be spared from a delete-by-number list.
    """
    path = os.path.join(unpacked_dir, "ppt", "slides", slide_filename)
    if not os.path.exists(path):
        return ""
    try:
        with open(path, encoding="utf-8") as fh:
            xml = fh.read()
    except OSError:
        return ""
    return " ".join(re.findall(r"<a:t>(.*?)</a:t>", xml, re.S))


def _set_run_text(paragraph, index, text):
    runs = paragraph.runs
    if index < len(runs):
        runs[index].text = text


def _rebuild_narrative(text_frame, segments, font_pt: float | None = None):
    """Rewrites a (possibly multi-run, multi-paragraph) narration text box as
    ONE paragraph made of `segments` = [(text, bold), ...], reusing the
    template's own bold/plain run styles so the look matches the original.

    `font_pt` overrides the template's size for every run — used by
    `_narration` to fit the text to its band. Colour, face and the bold/plain
    distinction are still taken from the template.
    """
    paragraphs = text_frame.paragraphs
    bold_style = plain_style = None
    for p in paragraphs:
        for r in p.runs:
            if r.font.bold and bold_style is None:
                bold_style = r.font
            if not r.font.bold and r.text.strip() and plain_style is None:
                plain_style = r.font
    if bold_style is None and paragraphs and paragraphs[0].runs:
        bold_style = paragraphs[0].runs[0].font
    if plain_style is None:
        plain_style = bold_style

    first_p = paragraphs[0]
    # clear all runs in the first paragraph
    for r in list(first_p.runs):
        r._r.getparent().remove(r._r)
    # drop any extra paragraphs
    for p in list(paragraphs[1:]):
        p._p.getparent().remove(p._p)

    for text, bold in segments:
        run = first_p.add_run()
        run.text = text
        src = bold_style if bold else plain_style
        if src is not None:
            run.font.bold = bold
            if src.size is not None:
                run.font.size = src.size
            if src.name:
                run.font.name = src.name
            try:
                if src.color and src.color.type is not None:
                    run.font.color.rgb = src.color.rgb
            except Exception:
                pass
        if font_pt is not None:
            run.font.size = Pt(font_pt)


def _find_narration(slide):
    """The narration text box on a slide cloned from the Team Deployment layout.

    Resolved by ROLE, with the sample template's wording only as a first
    preference. Every one of these slides used to be found by searching for the
    literal words "were deployed" — the phrase in the sample template's
    deployment sentence. On a template worded differently that matches nothing:
    `_narration` is never called, so the box keeps the template's own geometry
    (full width, 18pt, no height limit) while the chart and table are placed at
    the pipeline's coordinates, and the two collide. No error is raised and
    nothing is logged, so it looks like the layout fix simply did not apply.

    Fallback is the largest text shape that is not the title and sits above the
    content band — on this layout that is the narration and nothing else.
    """
    title = None
    try:
        title = slide.shapes.title
    except (AttributeError, KeyError):
        pass

    candidates = [sh for sh in slide.shapes
                  if sh.has_text_frame and sh is not title
                  and sh.text_frame.text.strip()]
    for sh in candidates:
        if "were deployed" in sh.text_frame.text.lower():
            return sh

    above = [sh for sh in candidates
             if sh.top is not None and sh.top < _in(CONTENT_TOP_IN)]
    if above:
        best = max(above, key=lambda sh: len(sh.text_frame.text))
        print(f"  narration: template wording not found, using {best.name!r}")
        return best
    print("  WARNING: no narration shape found on this slide — its text was "
          "left as the template had it")
    return None


def _find_landing_title(slide):
    """The title-bearing shape on the landing page, or None.

    Resolved by ROLE, then by size — never by the words already in it. The
    previous version matched on the text containing "Tracking Report", which
    works on the sample template and silently does nothing on any org template
    whose title slide happens to be worded differently: no error, no warning,
    the old title just stays. Whichever route succeeds is printed, so a run
    that did not rewrite the title says so.

    Order:
      1. the real title placeholder (CENTER_TITLE / TITLE)
      2. any shape whose text still mentions a tracking report
      3. the largest-typed text shape in the upper two-thirds of the slide
    """
    try:
        if slide.shapes.title is not None:
            print("  landing title: using the slide's title placeholder")
            return slide.shapes.title
    except (AttributeError, KeyError):
        pass

    texts = [sh for sh in slide.shapes
             if sh.has_text_frame and sh.text_frame.text.strip()]
    for sh in texts:
        if "tracking report" in sh.text_frame.text.lower():
            print(f"  landing title: matched on text in {sh.name!r}")
            return sh

    def biggest_pt(sh):
        sizes = [r.font.size.pt for p in sh.text_frame.paragraphs
                 for r in p.runs if r.font.size is not None]
        return max(sizes) if sizes else 0.0

    upper = [sh for sh in texts
             if sh.top is not None and sh.top < _in(SLIDE_H_IN * 0.66)]
    if upper:
        best = max(upper, key=lambda sh: (biggest_pt(sh), sh.width or 0))
        print(f"  landing title: falling back to the largest text shape "
              f"{best.name!r} ({biggest_pt(best):g}pt)")
        return best
    return None


def _set_landing_title(shape, state: str) -> None:
    """Write the three-line landing-page title into the template's own title box.

    Each line is a separate paragraph cloned from the template's, so the
    centring, typeface and colour are the template's own and only the size and
    weight are set here. The frame is centred vertically as well, so a
    three-line block sits balanced in the box instead of hanging off the bottom
    edge the way the template's single line was anchored to.
    """
    tf = shape.text_frame
    tf.word_wrap = True
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass

    first = tf.paragraphs[0]
    for extra in list(tf.paragraphs[1:]):
        extra._p.getparent().remove(extra._p)

    lines = [TITLE_LINE_1.format(state=state), TITLE_LINE_2, TITLE_LINE_3]
    targets = [first]
    for _ in lines[1:]:
        clone = copy.deepcopy(first._p)
        targets[-1]._p.addnext(clone)
        targets.append(_Paragraph(clone, tf))

    for para, text, size, bold in zip(targets, lines, TITLE_LINE_SIZES, TITLE_LINE_BOLD):
        runs = para.runs
        if runs:
            runs[0].text = text
            for extra in runs[1:]:
                extra._r.getparent().remove(extra._r)
            run = runs[0]
        else:
            run = para.add_run()
            run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        para.alignment = PP_ALIGN.CENTER


def _fit_font_pt(text: str, width_in: float, height_in: float,
                 base: float = NARR_BASE_PT, minimum: float = NARR_MIN_PT) -> float:
    """Largest font size from `base` down to `minimum` that fits `text` in a box.

    A character-count estimate, not a text metric — python-pptx cannot measure
    rendered text, and the deck is a Google Slides export whose font is usually
    substituted at render time anyway. The ratios above are set pessimistically
    so the estimate errs towards a slightly smaller narration rather than one
    that reaches into the chart.
    """
    n = max(len(text.strip()), 1)
    pt = base
    while pt > minimum:
        chars_per_line = max(1, int(width_in * 72.0 / (pt * NARR_CHAR_W_RATIO)))
        lines = math.ceil(n / chars_per_line)
        if lines * pt * NARR_LINE_H_RATIO / 72.0 <= height_in:
            return pt
        pt -= 0.5
    return minimum


def _narration(shape, segments, content_top_in: float = CONTENT_TOP_IN,
               top_in: float = NARR_TOP_IN, left_in: float = NARR_LEFT_IN,
               width_in: float = NARR_WIDTH_IN,
               max_pt: float = NARR_BASE_PT) -> float:
    """Write a narration into the fixed band above `content_top_in`.

    Pins the box to the band, turns off any autofit that would let it grow,
    and picks the largest font that still fits. Returns the size used, so the
    run log can show when a slide is being asked to carry too much text —
    hitting NARR_MIN_PT is the signal to shorten the wording, not to widen the
    box into the charts.
    """
    band_h = max(content_top_in - top_in - NARR_GAP_IN, 0.30)
    shape.left, shape.top = _in(left_in), _in(top_in)
    shape.width, shape.height = _in(width_in), _in(band_h)

    tf = shape.text_frame
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.TOP
    except Exception:
        pass
    for attr in ("margin_top", "margin_bottom"):
        try:
            setattr(tf, attr, Emu(0))
        except Exception:
            pass

    text = "".join(t for t, _ in segments)
    usable_w = width_in - 0.20
    pt = _fit_font_pt(text, usable_w, band_h, base=max_pt)
    _rebuild_narrative(tf, segments, font_pt=pt)

    # Where the text actually ends, by the same estimate used to size it. The
    # caller places the content below THIS rather than at a fixed constant, so
    # a narration that needs an extra line pushes the charts down instead of
    # running into them.
    chars_per_line = max(1, int(usable_w * 72.0 / (pt * NARR_CHAR_W_RATIO)))
    lines = math.ceil(max(len(text.strip()), 1) / chars_per_line)
    bottom_in = top_in + lines * pt * NARR_LINE_H_RATIO / 72.0
    return pt, bottom_in


def _fit_picture(slide, image_path, box):
    """Adds `image_path` to `slide`, contain-fit (not stretched) and centered
    within `box=(left, top, width, height)` in EMU — used both to replace an
    existing picture placeholder and to add a brand-new picture (e.g. the
    second, side-by-side chart on the merged cumulative-coverage slide)."""
    from PIL import Image
    box_left, box_top, box_w, box_h = box
    with Image.open(image_path) as im:
        img_w, img_h = im.size
    img_ratio = img_w / img_h
    box_ratio = box_w / box_h

    if img_ratio > box_ratio:
        new_w = box_w
        new_h = int(box_w / img_ratio)
    else:
        new_h = box_h
        new_w = int(box_h * img_ratio)
    left = box_left + (box_w - new_w) // 2
    top = box_top + (box_h - new_h) // 2
    return slide.shapes.add_picture(image_path, left, top, new_w, new_h)


# Title autofit. The template's master sets titles at 28pt, which suits its own
# short titles ("Background", "Team Deployed/Team Reported") but wraps to a
# second line on the longer ones this pipeline generates — and the second line
# lands on the accent bar underneath. Rather than shortening the wording, the
# font is stepped down just enough to keep one line, which preserves both the
# text and the template's look.
TITLE_BASE_PT = 28.0
TITLE_MIN_PT = 18.0
# Width of one character as a fraction of font size, for the template's bold
# sans title face. Measured at ~0.44 when that font is installed, but the deck
# is a Google Slides export and its font is often substituted — the substitute
# measured ~0.54 here. The higher figure is used because over-estimating only
# shrinks the title slightly more than needed, while under-estimating puts it
# back on two lines on whichever machine lacks the font.
TITLE_CHAR_WIDTH_RATIO = 0.56
# Placeholder is 7.52in wide with 0.1in insets each side.
TITLE_USABLE_WIDTH_IN = 7.25


def _set_title(slide, text: str) -> None:
    """Set a slide's title, shrinking the font if it would wrap to two lines."""
    title = slide.shapes.title
    if title is None:
        return
    paragraphs = title.text_frame.paragraphs
    if paragraphs and paragraphs[0].runs:
        for extra in list(paragraphs[0].runs[1:]):
            extra._r.getparent().remove(extra._r)
        run = paragraphs[0].runs[0]
    else:
        run = paragraphs[0].add_run()
    run.text = text

    est_in = len(text) * TITLE_CHAR_WIDTH_RATIO * TITLE_BASE_PT / 72.0
    if est_in > TITLE_USABLE_WIDTH_IN:
        fitted = TITLE_USABLE_WIDTH_IN * 72.0 / (len(text) * TITLE_CHAR_WIDTH_RATIO)
        run.font.size = Pt(max(TITLE_MIN_PT, round(fitted, 1)))


def _add_dataframe_table(slide, df, box, max_rows=None, font_size=10,
                         header_fill="1F4E79", overflow_label="more rows",
                         row_height=None, cell_margins=None, col_widths=None):
    """Render a small DataFrame as a native PowerPoint table inside `box`.

    A real table rather than a picture of one, so the numbers stay selectable
    and the deck keeps working if someone edits a figure before presenting.

    Long frames are truncated to `max_rows` with a trailing "+N more" row —
    the slide is a summary and the workbook carries the full list, but a
    silently cut table would be read as complete.

    `col_widths` is a tuple of INCHES, one per column. Without it PowerPoint
    distributes the width evenly, which gives a long LGA name the same room as
    a two-character hours figure; the name then wraps, the row grows to fit,
    and a table sized for N rows silently renders taller than its box and runs
    over whatever was placed beneath it.
    """
    from pptx.util import Pt as _Pt
    from pptx.dml.color import RGBColor as _RGB

    left, top, width, height = box
    body = df if max_rows is None else df.head(max_rows)
    hidden = len(df) - len(body)
    n_rows = len(body) + 1 + (1 if hidden > 0 else 0)
    n_cols = len(df.columns)
    if n_cols == 0 or n_rows <= 1:
        return None

    # Size the table to its content rather than stretching it to the box.
    # `row_height` caps it explicitly where the caller needs a compressed table
    # (the Time Spent slide, fitting ten LGAs); otherwise 0.30in a row.
    cap = int(row_height) if row_height else 274320
    row_h = Emu(int(min(int(height) / n_rows, cap)))
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width,
                                   Emu(int(row_h) * n_rows))
    table = shape.table
    # add_table distributes height evenly; set each row so the cap is honoured
    for r in table.rows:
        r.height = row_h
    if col_widths and len(col_widths) == n_cols:
        for c, w_in in enumerate(col_widths):
            table.columns[c].width = _in(w_in)
    if cell_margins:
        # a row cannot render shorter than its text plus the cell's margins,
        # so trimming the margins is what actually lets a compressed table be
        # compressed — setting the row height alone does nothing
        vert, horiz = cell_margins
        for r in range(n_rows):
            for c in range(n_cols):
                cell = table.cell(r, c)
                cell.margin_top = vert
                cell.margin_bottom = vert
                cell.margin_left = horiz
                cell.margin_right = horiz

    for c, col in enumerate(df.columns):
        cell = table.cell(0, c)
        cell.text = str(col)
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = _Pt(font_size)
                r.font.bold = True
                r.font.color.rgb = _RGB(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _RGB.from_string(header_fill)

    for i, (_, row) in enumerate(body.iterrows(), start=1):
        for c, value in enumerate(row):
            cell = table.cell(i, c)
            cell.text = "" if pd.isna(value) else (
                f"{value:,}" if isinstance(value, (int, float)) else str(value))
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = _Pt(font_size)

    if hidden > 0:
        cell = table.cell(n_rows - 1, 0)
        label = overflow_label if hidden != 1 else re.sub(r"s\b", "", overflow_label)
        cell.text = f"+{hidden:,} {label}"
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = _Pt(font_size - 1)
                r.font.italic = True
    return shape


def _table_caption(slide, table_shape, box, text: str) -> None:
    """Italic caption directly beneath a table, clear of its last row.

    Positioned from the RENDERED row height, not the shape's reported height —
    the latter is only what was requested, and the renderer grows rows to fit
    their text, which would print the caption across the table's last rows.
    """
    if table_shape is None:
        return
    n_rows = len(table_shape.table.rows)
    rendered = int(n_rows * TABLE_RENDERED_ROW_IN * 914400)
    top = int(table_shape.top) + max(int(table_shape.height), rendered) + int(_in(0.06))
    height = _in(0.36)
    # never let the caption run off the bottom of the slide
    top = min(top, int(_in(CONTENT_BOTTOM_IN)) - int(height))
    cap = slide.shapes.add_textbox(box[0], Emu(top), box[2], height)
    cap.text_frame.word_wrap = True
    try:
        cap.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    run = cap.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(TABLE_CAPTION_PT)
    run.font.italic = True


def _placeholder_note(slide, box, text: str, font_pt: float = 11) -> None:
    """Stand-in for a table that has no rows — an empty column reads as a
    rendering failure, while "nothing to follow up" is a finding in itself."""
    tb = slide.shapes.add_textbox(box[0], box[1], box[2], _in(0.80))
    tb.text_frame.word_wrap = True
    try:
        tb.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(font_pt)
    run.font.bold = True


def _overlap_report(prs, skip_slides=(1, 2)) -> list[str]:
    """Names any pair of visible shapes whose boxes intersect meaningfully.

    A build-time check, not a fixer. Slide geometry here is explicit, so an
    overlap means a constant is wrong or a table rendered taller than planned —
    both worth seeing in the run log rather than discovering in the meeting.

    Scoped to what this module lays out, so the signal stays usable:

    - `skip_slides` drops the title and background slides, whose overlapping
      decorative shapes are the template's own design and not ours to change.
    - The title placeholder is excluded. Its box is 0.92in tall for text that
      renders at roughly half that, so it "overlaps" whatever sits below on
      every slide in the deck while looking perfectly clean.
    - Only pairs sharing more than MIN_OVERLAP_IN^2 are reported; smaller
      intersections are rounding, not collisions.
    """
    MIN_OVERLAP_IN2 = 0.06
    issues = []
    for idx, slide in enumerate(prs.slides, start=1):
        if idx in skip_slides:
            continue
        title_shape = slide.shapes.title
        boxes = []
        for sh in slide.shapes:
            if title_shape is not None and sh is title_shape:
                continue
            try:
                l, t, w, h = int(sh.left), int(sh.top), int(sh.width), int(sh.height)
            except (TypeError, ValueError):
                continue
            if w <= 0 or h <= 0:
                continue
            # background/decorative fills span most of the slide by design
            if w > int(_in(9.6)) and h > int(_in(5.0)):
                continue
            label = sh.name
            if sh.has_text_frame and sh.text_frame.text.strip():
                label = f"{sh.name} [{sh.text_frame.text.strip()[:28]}]"
            boxes.append((l, t, l + w, t + h, label))
        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                a, b = boxes[i], boxes[j]
                ow = min(a[2], b[2]) - max(a[0], b[0])
                oh = min(a[3], b[3]) - max(a[1], b[1])
                if ow <= 0 or oh <= 0:
                    continue
                area_in2 = (ow / 914400.0) * (oh / 914400.0)
                if area_in2 > MIN_OVERLAP_IN2:
                    issues.append(f"slide {idx}: {a[4]} overlaps {b[4]} "
                                  f"({area_in2:.2f} sq in)")
    return issues


def _replace_picture(slide, shape, image_path, box=None):
    """Swaps a picture shape for a new image, fitting it within the original
    shape's bounding box (contain-fit, centered) rather than stretching it —
    our chart/map images rarely share the template screenshot's aspect ratio.
    Pass `box=(left, top, width, height)` in EMU to use a different box than
    the shape's own (the per-LGA template slide's box was hand-sized for one
    specific LGA's screenshot and is too small/oddly-shaped for the rest)."""
    box = box if box else (shape.left, shape.top, shape.width, shape.height)
    shape._element.getparent().remove(shape._element)
    return _fit_picture(slide, image_path, box)


def _pictures_by_x(slide):
    return sorted([s for s in slide.shapes if s.shape_type == 13],
                  key=lambda s: s.left)


def _delete_shape(shape):
    shape._element.getparent().remove(shape._element)


def _set_single_placeholder_bullet(shape, text):
    tf = shape.text_frame
    paragraphs = tf.paragraphs
    first_p = paragraphs[0]
    for r in list(first_p.runs):
        r._r.getparent().remove(r._r)
    for p in list(paragraphs[1:]):
        p._p.getparent().remove(p._p)
    run = first_p.add_run()
    run.text = text


def _lga_maps_are_side_by_side() -> bool:
    """Whether stage 4 rendered the LGA maps as two panels on one page.

    Read back from the layout stage 4 resolved rather than tracked through the
    call chain: this runs in the same process right after the maps were built,
    against the same template cache and the same template directories, so it
    reports what was actually drawn. Any failure answers "no", which only costs
    a slightly less specific slide title.
    """
    try:
        from stage4_maps import template_for
        spec = template_for("lga")
        return spec is not None and spec.is_multi_map
    except Exception:
        return False


# ----------------------------------------------------------- data helpers
def _norm_lga(series):
    return (series.astype(str).str.replace("_", " ", regex=False)
            .str.replace(r"\s+", " ", regex=True).str.strip().str.title())


def _lga_ward_team_settlement_counts(dip, deploy_csv, n_lga_authoritative=None):
    lga_col = find_col(dip, "lga")
    ward_col = find_col(dip, "ward")
    lat_col = find_col(dip, "latitude")
    lon_col = find_col(dip, "longitude")

    # Some settlement lists have naming variants for the same LGA (e.g.
    # "Talata Mafara" vs "Talata_Mafara") that inflate a naive unique-count.
    # The map slides are generated per BOUNDARY polygon (the authoritative
    # source), so when that count is available it takes precedence — keeping
    # "N LGAs" in the narrative consistent with the actual number of LGA map
    # slides in the deck.
    n_lga = n_lga_authoritative if n_lga_authoritative is not None else _norm_lga(dip[lga_col].dropna()).nunique()
    n_ward = int(dip.groupby(_norm_lga(dip[lga_col]))[ward_col]
                 .nunique().sum()) if ward_col else 0
    n_geo = int(dip[[lat_col, lon_col]].notna().all(axis=1).sum()) if lat_col and lon_col else 0
    n_total = len(dip)

    n_teams = 0
    if deploy_csv and os.path.exists(deploy_csv):
        dep = pd.read_csv(deploy_csv)
        gt = dep[dep["LGA"] == "Grand Total"]
        if len(gt):
            n_teams = int(gt.iloc[0]["Teams Deployed"])
    return n_lga, n_ward, n_teams, n_geo, n_total


def _coverage_stats(dip, status_col):
    """Stats where `status_col` itself holds Visited/Not Visited text (cum_col).

    For the cumulative column this already reflects evidence of tracks from any
    day of the period, whichever day the settlement was planned for: stage 2
    marks `day_{N}_cumm` "Visited" on cumulative track evidence, so the rule is
    applied once at the source rather than re-derived by each consumer.
    """
    mask = dip[status_col].notna()
    visited = int((dip.loc[mask, status_col] == "Visited").sum())
    planned = int(mask.sum())
    pct = visited / planned * 100 if planned else 0
    return visited, planned, pct


def _daily_coverage_stats(dip, day, cum_col):
    """DAILY coverage — the reporting day alone, never Day 1 to date.

    BOTH sides of the ratio are scoped to the day:

    denominator  the settlements PLANNED for this day, read from the settlement
                 list's own day-N planning flag. Measuring the day's work
                 against the whole campaign list is what made a normal day look
                 like a 15% day — the teams were never sent to the other 85%.
                 Falls back to the full planned list when the settlement list
                 carries no such flag, and says so.
    numerator    `day_{N}_daily`, which stage 2 derives from THIS run's tracks
                 only, so no previous day's visit is inside it.

    Falls back to the cumulative column ONLY on a settlement list produced
    before stage 2 wrote the daily column, and warns when it does — a silent
    fallback would report Day 1's work as the current day's, which is the exact
    mistake the daily/cumulative split exists to prevent.
    """
    daily_col = f"day_{day}_daily"
    scheduled, scope_col = day_scope_mask(dip, day)
    mask = dip[cum_col].notna() & scheduled
    if scope_col:
        print(f"  daily coverage scoped to settlements planned for day {day} "
              f"via '{scope_col}' ({int(mask.sum()):,} settlements)")
    else:
        print(f"  no day-{day} planning flag in the settlement list — daily "
              f"coverage measured against the full planned list")

    planned = int(mask.sum())
    if daily_col in dip.columns:
        visited = int((dip.loc[mask, daily_col] == "Visited").sum())
        return visited, planned, (visited / planned * 100 if planned else 0), True
    print(f"  WARNING: no '{daily_col}' column — the daily slide will show "
          f"cumulative figures. Re-run stage 2 for true daily numbers.")
    visited = int((dip.loc[mask, cum_col] == "Visited").sum())
    return visited, planned, (visited / planned * 100 if planned else 0), False


# ------------------------------------------------------------------ build
def build_pptx_report(template_path: str, state: str, day: int, visitation_csv: str,
                      cum_col: str, deploy_csv: str, time_csv: str, maps_folder: str,
                      charts_folder: str, output_file: str,
                      campaign_name: str = "Vaccination Tracking Report",
                      state_logo_path: str | None = None,
                      tracks_file: str | None = None) -> str:
    """`tracks_file` is the merged track export. The Time Spent Analysis slide
    is measured from it directly (per team, per tracking date, from the
    timestamps); without it that slide's figures and table are omitted rather
    than filled from a different measure."""
    dip = pd.read_csv(visitation_csv, low_memory=False)
    lga_col = find_col(dip, "lga")
    daily_col = f"day_{day}_daily"

    lga_pngs = sorted(f for f in os.listdir(maps_folder)
                      if f.endswith(".png") and "statewide" not in f and "implementation" not in f)
    n_lga, n_ward, n_teams, n_geo, n_total = _lga_ward_team_settlement_counts(
        dip, deploy_csv, n_lga_authoritative=len(lga_pngs) or None)
    # Daily = this reporting day only. Cumulative = Day 1 to this day.
    day_visited, day_planned, day_pct, day_is_true_daily = _daily_coverage_stats(
        dip, day, cum_col)
    cum_visited, cum_planned, cum_pct = _coverage_stats(dip, cum_col)
    # counts behind the Cumulative Settlements Coverage slide — same helper the
    # chart uses, so slide narrative and chart can never disagree
    from stage_charts import cumulative_coverage_stats
    cum_cov = cumulative_coverage_stats(dip, cum_col)

    deploy_stats = {"deployed": 0, "reported": 0, "pending": 0, "reporting_pct": 0.0}
    if deploy_csv and os.path.exists(deploy_csv):
        dep = pd.read_csv(deploy_csv)
        gt = dep[dep["LGA"] == "Grand Total"]
        if len(gt):
            g = gt.iloc[0]
            deploy_stats = {"deployed": int(g["Teams Deployed"]), "reported": int(g["Teams Reported"]),
                           "pending": int(g["Teams Pending"]),
                           "reporting_pct": float(g["Reporting %"]) if "Reporting %" in g else 0.0}

    # "compliant" = teams that spent 1 hour or more in the field (categories
    # from stage3_erm_workbook.time_spent_analysis: "1 - 2 hrs" and ">2 hrs")
    time_stats = {"total": 0, "compliant": 0, "compliant_pct": 0.0, "no_tracks": 0}
    if time_csv and os.path.exists(time_csv):
        tdf = pd.read_csv(time_csv)
        gt = tdf[tdf["Time Spent"] == "Grand Total"]
        total = int(gt.iloc[0]["Number of Teams"]) if len(gt) else 0
        compliant = int(tdf.loc[tdf["Time Spent"].isin(["1 - 2 hrs", ">2 hrs"]), "Number of Teams"].sum())
        no_tracks_row = tdf[tdf["Time Spent"] == "0 (No evidence of tracks)"]
        time_stats = {
            "total": total, "compliant": compliant,
            "compliant_pct": compliant / total if total else 0.0,
            "no_tracks": int(no_tracks_row.iloc[0]["Number of Teams"]) if len(no_tracks_row) else 0,
        }

    # Daily time spent tracking — the ONLY measure on the Time Spent slide.
    # Measured from the track timestamps, per team, for one tracking date, so
    # the population is the teams that actually transmitted. Both the chart and
    # the LGA table beside it come from these rows, so the two cannot disagree.
    # Failure here costs the slide's figures, not the deck.
    daily_time_stats: dict = {}
    under_12_table = None
    if tracks_file and os.path.exists(tracks_file):
        try:
            from team_daily_time import (analyse as analyse_daily_time,
                                         scope_from_visitation, UNDER_12)
            # `dip` is stage 2's output and already covers only this campaign's
            # state, so it defines the scope. GTS exports are national; without
            # this the slide would count every state's teams.
            dts = analyse_daily_time(tracks_file, "latest",
                                     scope=scope_from_visitation(dip))
            daily_time_stats = dts["meta"]
            by_lga = dts.get("by_lga")
            if by_lga is not None and len(by_lga) and UNDER_12 in by_lga.columns:
                u12 = (by_lga[["LGA", UNDER_12]]
                       .rename(columns={UNDER_12: "Teams <12 min"}))
                u12 = u12[u12["Teams <12 min"] > 0].sort_values(
                    "Teams <12 min", ascending=False).reset_index(drop=True)
                if len(u12):
                    under_12_table = u12
        except Exception as exc:
            print(f"  daily time-spent detail unavailable for the slide ({exc})")
    else:
        print("  daily time-spent detail unavailable — no merged track export "
              "was passed to stage 6")

    # Estimated <5 target population and household coverage. The charts are
    # produced by stage 4; this only needs the figures for the narrative, and
    # whether there is anything to show at all — a settlement list with neither
    # estimate means the slide is not created rather than created empty.
    target_cov = None       # the DAILY analysis — what the slide leads with
    target_cov_cum = None   # cumulative, quoted after it as context
    try:
        from target_population import analyse_daily_and_cumulative
        both = analyse_daily_and_cumulative(visitation_csv, daily_col, cum_col)
        primary = both["daily"] or both["cumulative"]
        if primary and (primary["children"] or primary["households"]):
            target_cov = primary
            target_cov_cum = both["cumulative"] if both["daily"] else None
    except Exception as exc:
        print(f"  target population coverage unavailable for the slide ({exc})")

    # Team Time Efficiency — field minutes against grid-weighted coverage.
    # Optional like the rest: an older visitation CSV without the grid-cell
    # columns costs this slide, not the deck.
    eff = None
    try:
        from team_efficiency import analyse as analyse_team_efficiency
        eff = analyse_team_efficiency(visitation_csv, "cumulative")
        if not len(eff["per_team"]):
            eff = None
    except Exception as exc:
        print(f"  team efficiency unavailable for the slide ({exc})")

    # DAILY charts for slide 5. Both read the day-only columns stage 2 writes,
    # so nothing on that slide carries a previous day's work. The stacked bar
    # comes from `lga_range_stacked`, the same function that draws slide 6's —
    # identical bands, order, colours and labelling, only the scope differs.
    from stage_charts import state_coverage_donut, lga_range_stacked
    daily_donut = os.path.join(charts_folder, "state_coverage_donut_daily.png")
    daily_lga_bar = os.path.join(charts_folder, "lga_coverage_stacked_daily.png")
    state_coverage_donut(visitation_csv, cum_col, day, daily_donut,
                         label=f"State Daily Settlement Coverage — Day {day} only",
                         class_col="Daily Settlement Coverage", day_scope=True)
    lga_range_stacked(visitation_csv, "Daily Coverage", daily_col, daily_lga_bar,
                      label=f"State Daily Coverage by LGA — Day {day} only",
                      day_scope=day)

    # ---------------------------------------------------------- work dir
    workdir = tempfile.mkdtemp(prefix="gts_pptx_")
    unpacked = os.path.join(workdir, "unpacked")
    safe_extract(template_path, unpacked)

    # ------------------------------------------------- structural edits
    delete_slide(unpacked, SLIDE_DAY_COMPARE_TABLE)
    # cumulative-by-LGA is no longer its own slide — its chart moves onto
    # the cumulative-coverage slide (SLIDE_CUMULATIVE_DONUT) alongside the
    # donut, side by side, matching the daily-coverage slide's layout
    delete_slide(unpacked, SLIDE_CUMULATIVE_LGA)
    preserved_slides = []
    for s in SLIDE_LGA_MAP_ORIGINALS_TO_DROP:
        text = _slide_text(unpacked, s).lower()
        if any(kw in text for kw in PRESERVE_SLIDE_KEYWORDS):
            preserved_slides.append(s)
            print(f"  keeping {s} — matched {PRESERVE_SLIDE_KEYWORDS}")
            continue
        delete_slide(unpacked, s)

    # Time Spent Analysis has no dedicated slide in the template — duplicate
    # the Team Deployment slide (title + narrative + one chart, the same
    # layout we need) and insert the copy right after it
    _time_spent_slide_name = duplicate_slide(unpacked, SLIDE_DEPLOY, after=SLIDE_DEPLOY)

    # Cumulative Settlements Coverage likewise has no slide of its own —
    # another copy of the Team Deployment layout, inserted after the coverage
    # slides so it reads as their summary. This pushes the statewide map and
    # every LGA map slide down one index (see the content-edit section).
    cumulative_slide = duplicate_slide(unpacked, SLIDE_DEPLOY, after=SLIDE_CUMULATIVE_DONUT)

    # Target Population & Household Coverage — a new slide, again cloned from
    # the Team Deployment layout so it inherits the template's title style,
    # background and branding rather than introducing a new look. It sits after
    # the cumulative-coverage summary and before the maps, so the deck reads:
    # settlements covered -> what that means in children and households -> where.
    # Created only when the settlement list actually carries the estimates;
    # when it does, everything from the statewide map onward shifts down one
    # more index.
    # Team Time Efficiency follows Time Spent Analysis — same subject, finer
    # question — so the deployment / time / efficiency trio reads in order.
    efficiency_slide = None
    if eff is not None:
        efficiency_slide = duplicate_slide(unpacked, SLIDE_DEPLOY,
                                           after=_time_spent_slide_name)

    target_slide = None
    if target_cov is not None:
        target_slide = duplicate_slide(unpacked, SLIDE_DEPLOY, after=cumulative_slide)

    lga_slide_files = [SLIDE_LGA_MAP_TEMPLATE]
    prev = SLIDE_LGA_MAP_TEMPLATE
    for _ in range(len(lga_pngs) - 1):
        prev = duplicate_slide(unpacked, SLIDE_LGA_MAP_TEMPLATE, after=prev)
        lga_slide_files.append(prev)

    clean_unused_slides(unpacked)
    intermediate = os.path.join(workdir, "intermediate.pptx")
    rezip(unpacked, intermediate)

    # ------------------------------------------------------- content edits
    prs = Presentation(intermediate)
    slides = list(prs.slides)  # order now matches sldIdLst == final presentation order

    def shapes_by_name_prefix(slide):
        return {s.name: s for s in slide.shapes}

    # slide 1 — landing page. `campaign_name` is deliberately NOT on this slide:
    # the agreed title structure is the three lines in TITLE_LINE_*, and the
    # campaign name is not one of them. It still appears on the .docx title
    # page built by stage 5.
    s = slides[0]
    title_shape = _find_landing_title(s)
    if title_shape is not None:
        _set_landing_title(title_shape, state)
    else:
        print("  WARNING: no title shape found on the landing page — the "
              "three-line title was NOT written. Run "
              "`python inspect_template_title.py <your-template.pptx>` to see "
              "what slide 1 actually contains.")

    for sh in s.shapes:
        if not sh.has_text_frame or sh is title_shape:
            continue
        t = sh.text_frame.text
        if re.match(r"^Day\s+\d+$", t.strip()):
            _set_run_text(sh.text_frame.paragraphs[0], 0, f"Day {day}")
        elif re.match(r"^[A-Za-z]+ \d{1,2}, \d{4}$", t.strip()):
            _set_run_text(sh.text_frame.paragraphs[0], 0,
                          date.today().strftime("%B %d, %Y"))
    if state_logo_path and os.path.exists(state_logo_path):
        pic = next((sh for sh in s.shapes if sh.shape_type == 13), None)
        if pic is not None:
            _replace_picture(s, pic, state_logo_path,
                             box=(pic.left, pic.top, pic.width, pic.height))

    # slide 2 — background
    s = slides[1]
    for sh in s.shapes:
        if sh.has_text_frame and "supporting the" in sh.text_frame.text:
            paras = sh.text_frame.paragraphs
            paras[0].runs[0].text = (f"eHealth Africa is supporting the {state} campaign with "
                                     f"vaccination tracking across the following:")
            body_lines = [
                f"{n_lga} LGAs", f"{n_ward} Wards",
                f"{n_teams:,} vaccination teams are being tracked",
                f"{n_geo:,} settlements with geo-coordinates across the state are being tracked",
                f"{n_total:,} total settlements with and without geo-coordinates planned",
            ]
            li = 0
            for p in paras[2:]:
                if p.runs and li < len(body_lines):
                    p.runs[0].text = body_lines[li]
                    li += 1
        elif sh.shape_type == 13:
            impl = os.path.join(maps_folder, f"{state}_implementation_map.png")
            if os.path.exists(impl):
                _replace_picture(s, sh, impl)

    # slide 3 — team deployed/reported
    s = slides[2]
    narr = _find_narration(s)
    if narr is not None:
        _narration(narr, [
                (f"{deploy_stats['reported']:,}", True),
                (" of ", False), (f"{deploy_stats['deployed']:,} ", True),
                ("deployed teams have reported tracks — ", False),
                (f"{deploy_stats['reporting_pct']:.1%}", True),
                (" reporting, ", False), (f"{deploy_stats['pending']:,} ", True),
                ("still pending at the time of this report.", False),
            ], content_top_in=1.75)
    for sh in s.shapes:
        if sh.shape_type == 13:
            dep_chart = os.path.join(charts_folder, "team_deployment.png")
            if os.path.exists(dep_chart):
                _replace_picture(s, sh, dep_chart)

    # slide 4 — Time Spent Analysis (duplicated from the Team Deployment slide
    # above, so it arrives as title + narrative + one full-width chart).
    #
    # ONE measure only: for each team, how long it was actually tracking on the
    # reporting date, read from the track timestamps (`team_daily_time`). The
    # settlement-derived range analysis and the 08:00-15:00 field window used
    # to appear here too; both counted a different population from the bars
    # beside them, so the slide could contradict itself. The chart on the left
    # and the LGA table on the right now come from the same rows.
    s = slides[3]
    _set_title(s, f"Time Spent Analysis - Day {day}")

    time_chart = os.path.join(charts_folder, "daily_time_spent.png")
    has_time_chart = os.path.exists(time_chart)
    narr = _find_narration(s)
    if narr is not None:
        if daily_time_stats.get("total_teams"):
            m = daily_time_stats
            segments = [
                (f"{m['total_teams']:,} ", True),
                ("teams transmitted tracks on ", False),
                (f"{m['date']}", True),
                (", with a median of ", False),
                (f"{m['median_minutes']:.0f} minutes ", True),
                ("tracking. ", False),
                (f"{m['teams_under_12']:,} ", True),
                (f"({m['under_12_pct']:.1%})", True),
                (" tracked for under 12 minutes — too little contact time to "
                 "have worked their settlements, and the follow-up list for "
                 "tomorrow's takeoff.", False),
            ]
        else:
            segments = [
                ("No team tracking time could be measured for this reporting "
                 "date — check that the track export covers it.", False),
            ]
        _narration(narr, segments)

    for sh in list(s.shapes):
        if sh.shape_type == 13:
            # the cloned layout's single chart placeholder becomes the left-hand
            # chart; the table is added beside it below
            if has_time_chart:
                _replace_picture(s, sh, time_chart, box=TIME_CHART_BOX)
            else:
                _delete_shape(sh)

    if has_time_chart and under_12_table is not None and len(under_12_table):
        table_shape = _add_dataframe_table(
            s, under_12_table, TIME_TABLE_BOX, max_rows=TIME_TABLE_MAX_ROWS,
            font_size=TIME_TABLE_FONT_PT, row_height=TIME_TABLE_ROW_HEIGHT,
            overflow_label="more LGAs", col_widths=TIME_TABLE_COL_WIDTHS,
            cell_margins=(TIME_TABLE_CELL_MARGIN_V, TIME_TABLE_CELL_MARGIN_H))
        _table_caption(s, table_shape, TIME_TABLE_BOX, TIME_TABLE_CAPTION)
    elif has_time_chart:
        _placeholder_note(
            s, TIME_TABLE_BOX,
            "No team tracked for under 12 minutes on this date.")

    # slide 5 (when present) — TEAM TIME EFFICIENCY. Sits directly after Time
    # Spent Analysis: same subject, finer question. Everything from the daily
    # coverage slide onward shifts down one index when it is created.
    eff_offset = 0
    if efficiency_slide is not None:
        eff_offset = 1
        s = slides[4]
        _set_title(s, EFF_SLIDE_TITLE)
        m = eff["meta"]
        eff_png = os.path.join(charts_folder, "team_efficiency.png")
        # The narration is the analytical finding, not a description of the
        # scatter. Two clauses at most: how many teams put in the hours without
        # the coverage to show for it, and — the sharper finding — how many
        # reported working while their tracks never left one place.
        segments = [
            (f"{m['flagged_teams']:,} of {m['total_teams']:,} teams "
             f"({m['flagged_pct']:.1%})", True),
            (" logged more than 2 hours in the field while covering less than "
             "50% of their assigned grid. This indicates substantial time "
             "spent on the ground without a corresponding level of coverage.",
             False),
        ]
        if m.get("stationary_only"):
            # "A further" means ADDITIONAL to the teams already counted above.
            # `stationary_teams` includes those that are also over 2 hours with
            # low coverage, so quoting it here would count them twice.
            segments += [
                (" A further ", False),
                (f"{m['stationary_only']:,} teams "
                 f"({m['stationary_only_pct']:.1%})", True),
                (" reported working but moved through no more than one grid "
                 "cell, with a median of ", False),
                (f"{m['stationary_median_hours']:.1f} hours", True),
                (" of recorded tracks from effectively one location. These "
                 "teams should be ", False),
                ("flagged for supervisory verification", True), (".", False),
            ]
        else:
            segments += [(" No team's tracks showed it stationary in one "
                          "location for the day.", False)]

        # Narration first, then the content placed BELOW where it actually
        # ended. A fixed content top only works while the narration fits the
        # band it was sized for; deriving the top from the narration means an
        # extra line pushes the chart and table down rather than colliding
        # with them, whatever the wording or the rendering font turns out to be.
        narr = _find_narration(s)
        eff_top = EFF_CONTENT_TOP_IN
        if narr is not None:
            _pt, narr_bottom = _narration(
                narr, segments, content_top_in=EFF_CONTENT_TOP_IN,
                max_pt=EFF_NARR_MAX_PT)
            eff_top = max(EFF_CONTENT_TOP_IN, narr_bottom + EFF_NARR_CLEARANCE_IN)
            if eff_top > EFF_CONTENT_TOP_IN:
                print(f"  efficiency slide: narration runs to "
                      f"{narr_bottom:.2f}in, content pushed to {eff_top:.2f}in")

        # The follow-up table used to sit to the right of the chart. It has
        # been removed, so the chart takes the full width of the slide — which
        # is what makes the flagged team codes in its right-hand gutter legible
        # at presentation size rather than cramped into half a slide.
        eff_chart_box = (_in(EFF_CHART_LEFT_IN), _in(eff_top),
                         _in(EFF_CHART_WIDTH_IN),
                         _in(max(1.2, CONTENT_BOTTOM_IN - eff_top)))

        for sh in list(s.shapes):
            if sh.shape_type == 13:
                if os.path.exists(eff_png):
                    _replace_picture(s, sh, eff_png, box=eff_chart_box)
                else:
                    _delete_shape(sh)

    # slide 6 — STATE DAILY COVERAGE: this reporting day only, no prior days
    s = slides[4 + eff_offset]
    for sh in list(s.shapes):
        if sh.has_text_frame:
            txt = sh.text_frame.text
            if txt.strip().startswith("State Daily Coverage"):
                _set_title(s, f"State Daily Coverage - Day {day}")
            elif "settlements planned for" in txt or ("of the" in txt and "planned for" in txt):
                _narration(sh, [
                    (f"{day_visited:,}", True), (" of the ", False),
                    (f"{day_planned:,} ", True),
                    ("geo-coded settlements scheduled for ", False),
                    (f"Day {day} ", True), ("were visited ", False),
                    (f"({day_pct:.1f}%)", True),
                    (", from this day's tracks alone across ", False),
                    (f"{n_lga} LGAs", True), (".", False),
                    ("" if day_is_true_daily else
                     " Daily status unavailable — figures shown are cumulative.", False),
                ], content_top_in=COVERAGE_TOP_IN)
            elif re.match(r"^[\d,]+\s*\([\d.]+%\)$", txt.strip()) or re.match(r"^[\d,]+$", txt.strip()):
                _delete_shape(sh)  # redundant floating number overlay — baked into our chart image
    pics = _pictures_by_x(s)
    daily_pic_boxes = COVERAGE_PIC_BOXES
    if len(pics) >= 2:
        _replace_picture(s, pics[0], daily_donut, box=daily_pic_boxes[0])       # left = donut
        pics2 = _pictures_by_x(s)
        _replace_picture(s, pics2[-1], daily_lga_bar, box=daily_pic_boxes[-1])  # right = LGA stacked bar
    elif len(pics) == 1:
        _replace_picture(s, pics[0], daily_donut, box=daily_pic_boxes[0])
        _fit_picture(s, daily_lga_bar, daily_pic_boxes[-1])

    # slide 6 — STATE CUMULATIVE COVERAGE: Day 1 through the reporting day.
    # Same two-chart layout as slide 5 and the same stacked-bar design, so the
    # two slides can be read against each other; only the scope differs.
    s = slides[5 + eff_offset]
    for sh in list(s.shapes):
        if sh.has_text_frame and (sh.text_frame.text.strip().startswith("Cumulative Settlements Coverage")
                                  or sh.text_frame.text.strip().startswith("Cumulative Coverage by LGA")):
            _set_title(s, f"State Cumulative Coverage - Day 1 to Day {day}")
        elif sh.has_text_frame and "cumulative" in sh.text_frame.text.lower():
            _narration(sh, [
                (f"{cum_visited:,}", True), (" of ", False),
                (f"{cum_planned:,} ", True),
                ("planned geo-coded settlements ", False),
                (f"({cum_pct:.1f}%)", True),
                (" have evidence of tracks from ", False),
                (f"Day 1 to Day {day}", True),
                (" across ", False), (f"{n_lga} LGAs", True),
                (", counted whichever day each was scheduled for. ", False),
                (f"{cum_planned - cum_visited:,} ", True),
                ("remain unreached.", False),
            ], content_top_in=1.66)
        elif sh.shape_type == 13:
            cum_donut = os.path.join(charts_folder, "state_coverage_donut.png")
            box = daily_pic_boxes[0]
            if os.path.exists(cum_donut):
                _replace_picture(s, sh, cum_donut, box=box)  # left = donut
    lga_cum = os.path.join(charts_folder, "lga_cumulative_counts.png")
    if os.path.exists(lga_cum):
        _fit_picture(s, lga_cum, daily_pic_boxes[-1])         # right = LGA cumulative bar

    # slide 7 — Cumulative Settlements Coverage (index 6: the second Team
    # Deployment duplicate, inserted after the cumulative-coverage slide).
    # Same layout as Time Spent: title + narrative + one chart.
    s = slides[6 + eff_offset]
    _set_title(s, "Cumulative Settlements Coverage")
    narr = _find_narration(s)
    if narr is not None:
            if cum_cov:
                _narration(narr, [
                    (f"{cum_cov['covered']:,} of {cum_cov['planned']:,} ", True),
                    ("planned settlements have been reached ", False),
                    (f"({cum_cov['covered_pct']:.1%})", True), (": ", False),
                    (f"{cum_cov['fully']:,} ", True), ("fully, ", False),
                    (f"{cum_cov['partial']:,} ", True), ("partially and ", False),
                    (f"{cum_cov['low']:,} ", True), ("at low coverage. ", False),
                    (f"{cum_cov['not_yet']:,} ", True),
                    (f"({cum_cov['not_yet_pct']:.1%})", True),
                    (" are still unvisited and carry into the remaining days.", False),
                ], content_top_in=1.75)
    for sh in s.shapes:
        if sh.shape_type == 13:
            cum_chart = os.path.join(charts_folder, "state_cumulative_coverage.png")
            if os.path.exists(cum_chart):
                _replace_picture(s, sh, cum_chart)

    # slide 8 (when present) — Target Population and Household Coverage.
    # Cloned from the Team Deployment layout, so the title style, background
    # and branding are the template's own; only the content band below the
    # narrative is rebuilt, into two estimate donuts plus a per-LGA bar.
    map_offset = 0
    if target_slide is not None:
        map_offset = 1
        s = slides[7 + eff_offset]
        _set_title(s, TARGET_SLIDE_TITLE)

        # DAILY finding first, then the cumulative position, then — only where
        # the settlement list supplies it — the estimated targeted missing
        # households. Kept to three short clauses: the charts below already
        # show the split, so the words carry the implication instead.
        kids, hh = target_cov["children"], target_cov["households"]
        miss = target_cov.get("missing_households")
        lead = kids or hh
        segments = [
            (f"{lead['within_reached']:,}", True),
            (" of the ", False), (f"{lead['visited_total']:,} ", True),
            ("estimated target population in settlements visited on ", False),
            (f"Day {day} ", True), ("were reached ", False),
            (f"({lead['within_pct']:.1%}", True), (" depth of coverage).", False),
        ]
        if target_cov_cum is not None:
            ck = target_cov_cum["children"] or target_cov_cum["households"]
            ck_pct = ck["within_reached"] / ck["total"] if ck["total"] else 0
            segments += [
                (" Cumulatively ", False), (f"{ck['within_reached']:,} ", True),
                ("of the ", False), (f"{ck['total']:,} ", True),
                ("campaign target has been reached ", False),
                (f"({ck_pct:.1%})", True), (".", False),
            ]
        if miss:
            segments += [
                (" An estimated ", False), (f"{miss['total']:,} ", True),
                ("targeted households are recorded as missing, ", False),
                (f"{miss['in_not_visited']:,} ", True),
                (f"({miss['in_not_visited_pct']:.0%})", True),
                (" of them in settlements not yet reached.", False),
            ]

        narr = _find_narration(s)
        if narr is not None:
            _narration(narr, segments)
        for sh in list(s.shapes):
            if sh.shape_type == 13:
                # the cloned layout's lone picture is replaced by the panels below
                _delete_shape(sh)

        panels = []
        if kids:
            panels.append(os.path.join(charts_folder, "target_children_donut.png"))
        if hh:
            panels.append(os.path.join(charts_folder, "target_households_donut.png"))
        lga_bar = os.path.join(charts_folder, "target_coverage_by_lga.png")
        if os.path.exists(lga_bar):
            panels.append(lga_bar)
        panels = [p for p in panels if os.path.exists(p)]
        boxes = TARGET_BOXES_3 if len(panels) >= 3 else TARGET_BOXES_2
        for png, box in zip(panels, boxes):
            _fit_picture(s, png, box)

    # statewide map — index 7, plus one if the target-population slide was
    # inserted ahead of it
    s = slides[7 + eff_offset + map_offset]
    for sh in s.shapes:
        if sh.has_text_frame:
            _set_title(s, f"State Coverage Map - Day {day}")
        elif sh.shape_type == 13:
            statewide = os.path.join(maps_folder, f"{state}_statewide_day_{day}.png")
            if os.path.exists(statewide):
                _replace_picture(s, sh, statewide, box=MAP_BOX)

    # slides 9.. — one per implementing LGA (names taken from the boundary
    # data itself, not reverse-parsed from filenames — safe_name() replaces
    # "/" with "-" for the filename, e.g. "Birnin Magaji/Kiyaw", so the
    # filename alone isn't a reliable source for the real LGA name)
    from stage4_maps import safe_name
    # A side-by-side map already carries "Visitation Status" / "Visitation
    # Coverage" titles over its two panels, so a slide titled "Coverage Map"
    # would describe only half of what is on screen.
    lga_slide_title = ("{lga} LGA Visitation Status & Coverage - Day {day}"
                       if _lga_maps_are_side_by_side()
                       else "{lga} LGA Coverage Map - Day {day}")
    lga_name_lookup = {safe_name(n): n for n in dip[lga_col].dropna().astype(str).str.strip().str.title().unique()}
    for i, png in enumerate(lga_pngs):
        s = slides[8 + eff_offset + map_offset + i]
        # strip the state prefix ONCE. An unbounded replace also removed the
        # state's name from the middle of the LGA's own name, so in Nasarawa
        # state "Nasarawa_Nasarawa_Egon_day_2.png" became "Egon" and
        # "Nasarawa_Nasarawa_day_2.png" became "day 2.png" — two slides titled
        # from the wrong string.
        stem = png.replace(f"{state}_", "", 1).replace(f"_day_{day}.png", "")
        lga_name = lga_name_lookup.get(stem, stem.replace("_", " "))
        for sh in s.shapes:
            if sh.has_text_frame:
                _set_title(s, lga_slide_title.format(lga=lga_name, day=day))
            elif sh.shape_type == 13:
                _replace_picture(s, sh, os.path.join(maps_folder, png), box=MAP_BOX)

    # challenges + photo slides — find by title text (positions shifted by now)
    for s in slides:
        if not s.shapes.title and not any(sh.has_text_frame for sh in s.shapes):
            continue
        title_text = ""
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                title_text = sh.text_frame.text.strip()
                break
        if title_text == "Challenges":
            for sh in s.shapes:
                if sh.has_text_frame and sh.text_frame.text.strip() != "Challenges":
                    _set_single_placeholder_bullet(sh, "[Add challenges observed today and recommendations.]")
                    break
        elif title_text in ("Deployment Pictures", "Supportive Supervision Pictures"):
            # The supervision slide now covers ERM as well as supportive
            # supervision, so it is retitled. Matched on the template's own
            # wording, not a slide index, since the index moves with the
            # number of LGA map slides.
            if title_text == "Supportive Supervision Pictures":
                _set_title(s, PICTURE_SLIDE_TITLES[title_text])
            for sh in list(s.shapes):
                if sh.shape_type == 13:
                    _delete_shape(sh)

    overlaps = _overlap_report(prs)
    if overlaps:
        print(f"  LAYOUT WARNING — {len(overlaps)} overlapping shape pair(s):")
        for line in overlaps[:12]:
            print(f"    {line}")
        if len(overlaps) > 12:
            print(f"    …and {len(overlaps) - 12} more")

    prs.save(output_file)
    shutil.rmtree(workdir, ignore_errors=True)
    print(f"Saved {output_file}")
    return output_file


if __name__ == "__main__":
    ap = argparse.ArgumentParser(description="Build the daily report in the org's PPTX template")
    ap.add_argument("--template", required=True)
    ap.add_argument("--state", required=True)
    ap.add_argument("--day", type=int, required=True)
    ap.add_argument("--visitation-csv", required=True)
    ap.add_argument("--cum-col", default=None)
    ap.add_argument("--deploy-csv", required=True)
    ap.add_argument("--time-csv", required=True)
    ap.add_argument("--maps-folder", required=True)
    ap.add_argument("--charts-folder", required=True)
    ap.add_argument("--output", required=True)
    ap.add_argument("--campaign-name", default="Vaccination Tracking Report")
    ap.add_argument("--state-logo", default=None, help="State PHC logo to place on the title slide")
    ap.add_argument("--tracks", default=None,
                    help="merged_tracks.csv — required for the Time Spent "
                         "Analysis slide, measured from the track timestamps")
    a = ap.parse_args()
    build_pptx_report(a.template, a.state, a.day, a.visitation_csv,
                      a.cum_col or f"day_{a.day}_cumm", a.deploy_csv, a.time_csv,
                      a.maps_folder, a.charts_folder, a.output, a.campaign_name,
                      state_logo_path=a.state_logo, tracks_file=a.tracks)
