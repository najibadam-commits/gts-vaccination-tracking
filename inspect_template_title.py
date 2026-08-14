"""Show what stage 6 will do to a template's landing page.

Answers "why didn't the title change?" without running the whole pipeline.
Prints every shape on slide 1, says which one stage 6 resolves as the title,
and — with --write — produces a one-slide proof file you can open.

    python inspect_template_title.py "organization's reporting template.pptx"
    python inspect_template_title.py my_template.pptx --state Kano --write

If the resolved title is "(none)", stage 6 cannot find a title shape on that
template and the landing page will be left untouched; send this output on and
the matcher can be taught the layout.
"""
import argparse
import os
import sys

from pptx import Presentation

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from stage6_pptx_report import (  # noqa: E402
    TITLE_LINE_1, TITLE_LINE_2, TITLE_LINE_3, TITLE_LINE_SIZES,
    TITLE_LINE_BOLD, _find_landing_title, _set_landing_title,
)


def describe(shape) -> str:
    try:
        ph = shape.placeholder_format.type
    except Exception:
        ph = None
    try:
        geom = (f"L={shape.left / 914400:.2f} T={shape.top / 914400:.2f} "
                f"W={shape.width / 914400:.2f} H={shape.height / 914400:.2f}")
    except Exception:
        geom = "geometry unavailable"
    sizes = []
    if shape.has_text_frame:
        sizes = [r.font.size.pt for p in shape.text_frame.paragraphs
                 for r in p.runs if r.font.size is not None]
    pt = f"{max(sizes):g}pt" if sizes else "inherited"
    text = shape.text_frame.text.replace("\n", " | ")[:70] if shape.has_text_frame else ""
    return (f"    id={shape.shape_id:<4} type={str(shape.shape_type):<18} "
            f"placeholder={str(ph):<16} {geom}  {pt}\n"
            f"      text: {text!r}")


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("template", help="the .pptx report template to inspect")
    ap.add_argument("--state", default="Nasarawa",
                    help="state name to preview in the title (default: Nasarawa)")
    ap.add_argument("--write", action="store_true",
                    help="save <template>_title_preview.pptx with the title applied")
    a = ap.parse_args()

    if not os.path.exists(a.template):
        print(f"No such file: {a.template}")
        return 2

    prs = Presentation(a.template)
    slide = list(prs.slides)[0]

    print(f"\nTemplate : {os.path.basename(a.template)}")
    print(f"Slides   : {len(prs.slides)}")
    print(f"Slide 1  : {len(slide.shapes)} shapes\n")
    for sh in slide.shapes:
        print(describe(sh))

    print("\n" + "-" * 68)
    print("Stage 6 resolves the landing-page title as:")
    target = _find_landing_title(slide)
    print(f"  -> {target.name!r} (id={target.shape_id})" if target is not None
          else "  -> (none)  <-- the landing page would be left unchanged")

    print("\nIt would be replaced by these three lines:")
    lines = [TITLE_LINE_1.format(state=a.state), TITLE_LINE_2, TITLE_LINE_3]
    for text, size, bold in zip(lines, TITLE_LINE_SIZES, TITLE_LINE_BOLD):
        print(f"  {size:>5g}pt {'bold  ' if bold else '      '}{text}")

    if a.write and target is not None:
        _set_landing_title(target, a.state)
        out = os.path.splitext(a.template)[0] + "_title_preview.pptx"
        prs.save(out)
        print(f"\nPreview saved: {out}")
    print()
    return 0 if target is not None else 1


if __name__ == "__main__":
    sys.exit(main())
